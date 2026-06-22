"""Document text extraction — turn office/PDF files into plain text.

A STRATEGY DISPATCH (suffix -> extractor function), shared by the `read` tool (reads a
document as text) and the Resource Manager describer (summarises a document). Adding a
format = add one extractor here; nothing else changes (Open/Closed). This is the *good*
kind of per-format code — each entry does real work, not a hardcoded label. Heavy deps
(python-docx, pypdf, openpyxl, python-pptx) are imported lazily, so a missing one only
disables that one format.
"""

from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger("agentd")
_MAX_EXTRACT_BYTES = 10_000_000   # don't text-extract very large documents


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    # data_only -> computed values (not formulas); read_only -> handles large sheets
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):  # skip fully-empty rows
                    parts.append(" | ".join(cells).rstrip(" |"))
        return "\n".join(parts)
    finally:
        wb.close()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


# suffix -> text extractor (documents are text-extracted instead of read as bytes)
EXTRACTORS = {
    ".docx": _extract_docx,
    ".pdf": _extract_pdf,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def is_document(path) -> bool:
    return Path(path).suffix.lower() in EXTRACTORS


def extract_text(path) -> str | None:
    """Extracted plain text for a supported document, else None (unsupported type,
    too large, missing dependency, or a corrupt file — never raises)."""
    p = Path(path)
    fn = EXTRACTORS.get(p.suffix.lower())
    if fn is None:
        return None
    try:
        if p.stat().st_size > _MAX_EXTRACT_BYTES:
            return None
        return fn(p)
    except Exception:  # noqa: BLE001 — missing dep / corrupt file -> just no extraction
        log.debug("document extract failed for %s", p, exc_info=True)
        return None
