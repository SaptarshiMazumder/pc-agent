"""export_pptx with per-object `image` elements and no base artwork: each object becomes its own
selectable PowerPoint picture (the apple/mango/banana-selectable ask, PPTX side)."""

import asyncio
import json
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

pytest.importorskip("pptx")
pytest.importorskip("PIL")

from PIL import Image  # noqa: E402


def test_export_pptx_objects_are_separate_pictures(tmp_path):
    from export_pptx_tool import ExportPptxTool

    Image.new("RGB", (80, 80), (200, 60, 60)).save(tmp_path / "o0.png")
    Image.new("RGB", (80, 80), (60, 60, 200)).save(tmp_path / "o1.png")
    spec = {
        "width": 400,
        "height": 300,
        "elements": [
            {"kind": "image", "path": str(tmp_path / "o0.png"), "x": 20, "y": 20, "width": 80, "height": 80},
            {"kind": "image", "path": str(tmp_path / "o1.png"), "x": 200, "y": 20, "width": 80, "height": 80},
            {"kind": "label", "text": "A", "x": 60, "y": 130},
        ],
    }
    jp = tmp_path / "spec.json"
    jp.write_text(json.dumps(spec), encoding="utf-8")
    out = tmp_path / "out.pptx"
    tool = ExportPptxTool(SimpleNamespace(workspace=str(tmp_path)))
    # NO artwork -> slide size from the spec; the image elements are the pictures
    res = asyncio.run(
        tool.execute("t", {"elements_path": str(jp), "out_path": str(out)}, None)
    )
    assert not res.is_error, res.content[0].text
    assert out.is_file()

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    prs = Presentation(str(out))
    slide = prs.slides[0]
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 2  # two SEPARATE, selectable object pictures (no slide-filling base)
    # slide sized from the spec (400x300 px @ 96dpi)
    assert abs(prs.slide_width - Emu(int(400 * 914400 / 96))) < Emu(10000)
    # at least one text box for the label
    texts = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    assert any("A" in t.text_frame.text for t in texts)
