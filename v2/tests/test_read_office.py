import asyncio
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "plugins" / "core_fs"))
from fs_tools import ReadTool  # built-in 'core_fs' bundle


async def _read(path: Path) -> str:
    res = await ReadTool(SimpleNamespace(workspace=path.parent)).execute(
        "c1", {"path": str(path)}, asyncio.Event()
    )
    assert res.is_error is False, res.content[0].text
    return res.content[0].text


@pytest.mark.asyncio
async def test_read_xlsx_extracts_cells(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Q1"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    p = tmp_path / "data.xlsx"
    wb.save(str(p))

    text = await _read(p)
    assert "# Sheet: Q1" in text
    assert "Name | Score" in text
    assert "Alice | 95" in text


@pytest.mark.asyncio
async def test_read_pptx_extracts_slides(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Ship Office reading"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    text = await _read(p)
    assert "# Slide 1" in text
    assert "Roadmap" in text and "Ship Office reading" in text


def test_read_tool_advertises_office_formats():
    assert ".xlsx" in ReadTool.description and ".pptx" in ReadTool.description
