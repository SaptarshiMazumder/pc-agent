"""export_pptx: build a fully-editable PowerPoint figure from the same render_editable_overlay-style `elements`
the figure pipeline already produces — the artwork as a picture, and EVERY annotation as a native,
restyleable PowerPoint object:

  • label       -> a text box, or a rounded-rectangle "pill" (fill/border) + centred text
  • annotation  -> a pill label + a leader connector + a target dot (the science callout)
  • arrow       -> a connector with a real arrowhead; multi-point paths become an elbow chain
  • leader      -> a thin connector + optional target dot
  • panel       -> a rounded-rectangle region frame + optional title
  • node        -> a flowchart node: rounded box + embedded icon PICTURE + wrapped text + step badge
  • dot         -> a small marker oval

The slide is sized to the artwork's pixels (96 dpi) so every element's pixel coordinate lands exactly
on the picture. Labels are retypable, boxes recolourable/resizable, connectors draggable — the "clean,
custom-styleable, resizable PPTX" deliverable.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from figexport_common import px, resolve_path

from agentd.application.interfaces.tool import Tool, ToolResult

# arrow `style` -> (stroke width pt, default colour) when the element omits width/color.
_ARROW_LOOK = {
    "plain": (2.0, "#374151"),
    "clean": (2.75, "#374151"),
    "biorender": (3.5, "#374151"),
    "subtle": (2.0, "#6b7280"),
    "flow": (2.75, "#374151"),
    "emphasis": (4.0, "#374151"),
    "activation": (2.75, "#2f855a"),
    "inhibition": (2.75, "#c53030"),
    "transport": (2.5, "#374151"),
}


class ExportPptxTool(Tool):
    name = "export_pptx"
    description = (
        "Export a fully-editable PowerPoint (.pptx) figure: the artwork as a picture, plus EVERY "
        "render_editable_overlay-style element as a native, restyleable object — label/annotation -> text box "
        "or rounded pill, arrow/leader -> connector with a real arrowhead (multi-point -> elbow chain), "
        "panel -> region frame, node -> flowchart box with an embedded ICON picture + wrapped text + "
        "step badge, dot -> marker. Slide is sized to the artwork's pixels so coordinates line up. "
        "Input: `artwork` (PNG/JPG), `out_path`, and `elements` inline OR `elements_path` (the spec "
        "file extract_annotations writes — never retype a large set). Everything is retypable / "
        "recolourable / resizable in PowerPoint."
    )
    label = "Export PPTX"
    concurrency = "parallel"
    parameters = {
        "type": "object",
        "required": ["artwork", "out_path"],
        "properties": {
            "artwork": {
                "type": "string",
                "description": "Path to the raster artwork (PNG/JPG). Use a transparent PNG if you want only the annotations.",
            },
            "out_path": {"type": "string", "description": "Output .pptx path."},
            "elements_path": {
                "type": "string",
                "description": "Path to a spec JSON — an array of elements or {width,height,elements} (what extract_annotations writes). Use instead of retyping a large `elements` list.",
            },
            "elements": {
                "type": "array",
                "description": "render_editable_overlay-style elements (label, annotation, arrow, leader, panel, node, dot). "
                "Coordinates are in the artwork's pixel space (origin top-left).",
                "items": {
                    "type": "object",
                    "required": ["kind"],
                    "properties": {"kind": {"type": "string"}},
                    "additionalProperties": True,
                },
            },
            "background": {
                "type": "string",
                "description": "Optional slide background hex (e.g. '#FFFFFF'). Default: none (white).",
            },
            "scale": {
                "type": "number",
                "description": "Optional size multiplier for text/strokes/dots (match render_editable_overlay's `scale`). Omit = AUTO from the artwork resolution so text isn't tiny on a 2K/4K slide.",
            },
        },
    }

    def __init__(self, config):
        self.config = config

    def _resolve(self, p):
        return resolve_path(self.config, p)

    def _elements(self, params: dict) -> list:
        """Elements inline OR from `elements_path` (array, or the {width,height,elements} spec
        extract_annotations writes)."""
        import json

        inline = params.get("elements")
        path = params.get("elements_path")
        if bool(inline) == bool(path):
            raise ValueError("provide exactly ONE of: elements, elements_path")
        if inline:
            return list(inline)
        data = json.loads(self._resolve(path).read_text(encoding="utf-8"))
        if isinstance(data, list):
            return data
        if isinstance(data, dict) and isinstance(data.get("elements"), list):
            return data["elements"]
        raise ValueError("`elements_path` must hold an array or {width,height,elements}")

    def _run(self, params: dict) -> dict:
        from PIL import Image
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.oxml.ns import qn
        from pptx.util import Emu, Pt

        E = lambda v: Emu(px(v))  # px -> EMU

        def rgb(c, default="#374151"):
            try:
                return RGBColor.from_string(str(c or default).lstrip("#"))
            except Exception:
                return RGBColor.from_string(default.lstrip("#"))

        def approx_w(text, fs, bold):
            k = 0.60 if bold else 0.54
            return max(24.0, len(str(text)) * fs * k)

        def set_text(shape, text, fs, color, bold=True, align="center", anchor="middle"):
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(px(6))
            tf.margin_top = tf.margin_bottom = Emu(px(3))
            try:
                tf.vertical_anchor = {
                    "top": MSO_ANCHOR.TOP,
                    "middle": MSO_ANCHOR.MIDDLE,
                    "bottom": MSO_ANCHOR.BOTTOM,
                }.get(anchor, MSO_ANCHOR.MIDDLE)
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.alignment = {
                "start": PP_ALIGN.LEFT,
                "left": PP_ALIGN.LEFT,
                "middle": PP_ALIGN.CENTER,
                "center": PP_ALIGN.CENTER,
                "end": PP_ALIGN.RIGHT,
                "right": PP_ALIGN.RIGHT,
            }.get(align, PP_ALIGN.CENTER)
            r = p.add_run()
            r.text = str(text)
            r.font.size = Pt(fs * 0.75)
            r.font.bold = bold
            r.font.color.rgb = rgb(color, "#1f2937")

        def rrect(x, y, w, h, fill, line, line_w=1.4, radius=None):
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(x), E(y), E(w), E(h))
            if fill is None:
                sp.fill.background()
            else:
                sp.fill.solid()
                sp.fill.fore_color.rgb = rgb(fill, "#ffffff")
            if line is None:
                sp.line.fill.background()
            else:
                sp.line.color.rgb = rgb(line, "#94a3b8")
                sp.line.width = Pt(line_w)
            if radius is not None:  # 0..1 corner roundness
                try:
                    sp.adjustments[0] = max(0.0, min(0.5, float(radius)))
                except Exception:
                    pass
            return sp

        def oval(x, y, d, fill):
            sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(x), E(y), E(d), E(d))
            sp.fill.solid()
            sp.fill.fore_color.rgb = rgb(fill, "#374151")
            sp.line.fill.background()
            return sp

        def _end(ln, tag, color, on):
            """Set a headEnd/tailEnd arrowhead on a connector/line's XML (python-pptx has no API)."""
            if not on:
                return
            try:
                el = ln._get_or_add_ln()
                el.append(
                    el.makeelement(qn(f"a:{tag}"), {"type": "triangle", "w": "med", "len": "med"})
                )
            except Exception:
                pass

        def _dash(ln, dash):
            if not dash:
                return
            try:
                el = ln._get_or_add_ln()
                el.append(el.makeelement(qn("a:prstDash"), {"val": "dash"}))
            except Exception:
                pass

        def connector(p0, p1, color, width=2.0, head_end=False, head_start=False, dash=False):
            c = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, E(p0[0]), E(p0[1]), E(p1[0]), E(p1[1])
            )
            c.line.color.rgb = rgb(color, "#374151")
            c.line.width = Pt(width)
            _end(c.line, "tailEnd", color, head_end)
            _end(c.line, "headEnd", color, head_start)
            _dash(c.line, dash)
            return c

        def polyline(points, color, width, head_end=True, head_start=False, dash=False):
            """Multi-point path -> a chain of straight connectors (fully editable elbow); arrowhead on
            the last segment only (and start head on the first)."""
            pts = [tuple(p) for p in points]
            if len(pts) < 2:
                return
            for i in range(len(pts) - 1):
                connector(
                    pts[i],
                    pts[i + 1],
                    color,
                    width,
                    head_end=(head_end and i == len(pts) - 2),
                    head_start=(head_start and i == 0),
                    dash=dash,
                )

        art = self._resolve(params["artwork"])
        with Image.open(art) as im:
            W, H = im.size
        # Resolution-aware sizing: mirror the SVG overlay engine so text/strokes are proportionate on a
        # 2K/4K slide instead of microscopic. Sizes are authored in ~1024px reference units; coordinates
        # stay in the artwork's pixel space.
        S = float(params["scale"]) if params.get("scale") else max(1.0, max(W, H) / 1024.0)
        prs = Presentation()
        prs.slide_width = E(W)
        prs.slide_height = E(H)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        if params.get("background"):
            bg = rrect(0, 0, W, H, params["background"], None, 0)  # noqa: F841  # TODO: bg is never added to the slide — the `background` param currently renders nothing
        slide.shapes.add_picture(str(art), Emu(0), Emu(0), width=E(W), height=E(H))

        n_text = n_shape = n_conn = n_pic = 0
        for el in self._elements(params):
            kind = el.get("kind")

            if kind == "label":
                text = el.get("text", "")
                x, y = float(el.get("x", 0)), float(el.get("y", 0))
                fs = float(el.get("font_size", 14)) * S
                bold = str(el.get("weight", "600")) in ("600", "700", "bold")
                anchor = el.get("anchor", "start")
                box = el.get("box")
                bw = approx_w(text, fs, bold) + 2 * 9 * S
                bh = fs + 2 * 6 * S
                left = x - (
                    bw / 2
                    if anchor in ("middle", "center")
                    else (bw if anchor in ("end", "right") else 0)
                )
                if box:
                    sp = rrect(
                        left,
                        y - bh / 2,
                        bw,
                        bh,
                        (box or {}).get("fill", "#ffffff"),
                        (box or {}).get("stroke", "#d1d5db"),
                        float((box or {}).get("stroke_width", 1)) * S,
                        radius=0.3,
                    )
                    set_text(sp, text, fs, el.get("color", "#1f2937"), bold, "center")
                    n_shape += 1
                else:
                    tb = slide.shapes.add_textbox(E(left), E(y - bh / 2), E(bw), E(bh))
                    set_text(
                        tb,
                        text,
                        fs,
                        el.get("color", "#1f2937"),
                        bold,
                        "left"
                        if anchor == "start"
                        else ("right" if anchor in ("end", "right") else "center"),
                    )
                    n_text += 1

            elif kind == "annotation":
                at = el.get("at", [0, 0])
                tgt = el.get("target", at)
                text = el.get("text", "")
                fs = float(el.get("font_size", 14)) * S
                bold = True
                bw = approx_w(text, fs, bold) + 2 * 9 * S
                bh = fs + 2 * 6 * S
                side = el.get("side", "left" if at[0] < tgt[0] else "right")
                edge = (at[0] + bw / 2, at[1]) if side == "left" else (at[0] - bw / 2, at[1])
                connector(
                    edge,
                    tgt,
                    el.get("leader_color", "#9ca3af"),
                    float(el.get("leader_width", 1.2)) * S,
                )
                if el.get("dot", True):
                    dr = float(el.get("dot_r", 3)) * S
                    oval(tgt[0] - dr, tgt[1] - dr, 2 * dr, el.get("dot_color", "#4b5563"))
                    n_shape += 1
                sp = rrect(
                    at[0] - bw / 2,
                    at[1] - bh / 2,
                    bw,
                    bh,
                    (el.get("box") or {}).get("fill", "#ffffff"),
                    (el.get("box") or {}).get("stroke", "#d1d5db"),
                    S,
                    radius=0.3,
                )
                set_text(sp, text, fs, el.get("color", "#1f2937"), bold, "center")
                n_shape += 1
                n_conn += 1

            elif kind == "arrow":
                pts = el.get("points") or [el.get("from", [0, 0]), el.get("to", [0, 0])]
                look = _ARROW_LOOK.get(el.get("style", "clean"), _ARROW_LOOK["clean"])
                width = float(el.get("width", look[0])) * S
                color = el.get("color", look[1])
                head = el.get("head", "standard") != "none"
                start_head = el.get("start_head", "none") not in ("none", None)
                polyline(
                    pts,
                    color,
                    width,
                    head_end=head,
                    head_start=start_head,
                    dash=bool(el.get("dash")),
                )
                n_conn += max(1, len(pts) - 1)

            elif kind == "leader":
                pts = el.get("points") or [el.get("from", [0, 0]), el.get("to", [0, 0])]
                color = el.get("color", "#6b7280")
                polyline(
                    pts,
                    color,
                    float(el.get("width", 1.2)) * S,
                    head_end=(el.get("head", "none") not in ("none", None)),
                    dash=bool(el.get("dash")),
                )
                if el.get("dot", el.get("head", "none") == "none"):
                    dr = float(el.get("dot_r", 2.6)) * S
                    oval(pts[-1][0] - dr, pts[-1][1] - dr, 2 * dr, el.get("dot_color", color))
                    n_shape += 1
                n_conn += max(1, len(pts) - 1)

            elif kind == "panel":
                x, y, w, h = float(el["x"]), float(el["y"]), float(el["w"]), float(el["h"])
                rrect(
                    x,
                    y,
                    w,
                    h,
                    el.get("fill"),
                    el.get("stroke", "#9ca3af"),
                    float(el.get("stroke_width", 1.4)) * S,
                    radius=0.04,
                )
                n_shape += 1
                if el.get("title"):
                    fs = float(el.get("title_size", 13)) * S
                    tw = approx_w(el["title"], fs, True) + 16 * S
                    tb = slide.shapes.add_textbox(
                        E(x + w / 2 - tw / 2), E(y - (fs + 6 * S) / 2), E(tw), E(fs + 8 * S)
                    )
                    set_text(tb, el["title"], fs, el.get("title_color", "#374151"), True, "center")
                    n_text += 1

            elif kind == "node":
                x, y, w, h = float(el["x"]), float(el["y"]), float(el["w"]), float(el["h"])
                pad = float(el.get("pad", 10)) * S
                sp = rrect(
                    x,
                    y,
                    w,
                    h,
                    el.get("fill", "#ffffff"),
                    el.get("stroke", "#94a3b8"),
                    float(el.get("stroke_width", 1.6)) * S,
                    radius=float(el.get("radius_frac", 0.1)),
                )
                n_shape += 1
                fs = float(el.get("font_size", 13)) * S
                text = el.get("text", "")
                icon = el.get("icon")
                if icon:
                    ip = self._resolve(icon)
                    if Path(ip).exists():
                        icon_h = (h - 2 * pad) * (0.58 if text else 1.0)
                        iw = min(w - 2 * pad, icon_h)
                        try:
                            slide.shapes.add_picture(
                                str(ip), E(x + (w - iw) / 2), E(y + pad), height=E(icon_h)
                            )
                            n_pic += 1
                        except Exception:
                            pass
                        if text:
                            tb = slide.shapes.add_textbox(
                                E(x + pad),
                                E(y + pad + icon_h),
                                E(w - 2 * pad),
                                E(h - 2 * pad - icon_h),
                            )
                            set_text(tb, text, fs, el.get("text_color", "#1f2937"), True, "center")
                            n_text += 1
                elif text:
                    set_text(sp, text, fs, el.get("text_color", "#1f2937"), True, "center")
                if el.get("step") is not None:
                    d = fs * 1.8
                    badge = oval(
                        x + 3 * S, y + 3 * S, d, el.get("step_bg", el.get("stroke", "#94a3b8"))
                    )
                    set_text(
                        badge,
                        str(el["step"]),
                        fs * 0.9,
                        el.get("step_color", "#ffffff"),
                        True,
                        "center",
                    )
                    n_shape += 1

            elif kind == "dot":
                dr = float(el.get("r", 3)) * S
                oval(float(el["x"]) - dr, float(el["y"]) - dr, 2 * dr, el.get("fill", "#374151"))
                n_shape += 1
            # 'raw' has no editable-PPTX equivalent -> skipped

        out = self._resolve(params["out_path"])
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        return {
            "out_path": str(out),
            "width": W,
            "height": H,
            "textboxes": n_text,
            "shapes": n_shape,
            "connectors": n_conn,
            "pictures": n_pic,
        }

    async def execute(self, tool_call_id, params, abort, on_update=None):
        try:
            r = await asyncio.to_thread(self._run, params)
        except Exception as e:
            return ToolResult.text(f"export_pptx failed: {e}", is_error=True)
        return ToolResult.text(
            f"PPTX -> {r['out_path']} (slide {r['width']}x{r['height']}px; {r['textboxes']} text box(es), "
            f"{r['shapes']} shape(s), {r['connectors']} connector(s), {r['pictures']} picture(s) — all editable).",
            details=r,
            artifacts=[r["out_path"]],
        )  # deliverable: the editable deck
