"""make_pptx: build a PowerPoint (.pptx) from a list of slides.

Each slide is a small spec — title, bullets or subtitle, an optional image (e.g. a diagram PNG),
and optional speaker notes (drop the narration text here). Uses python-pptx. 16:9 by default.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from agentd.application.interfaces.tool import Tool, ToolResult
from agentd.application.run_context import current_workspace


class MakePptxTool(Tool):
    name = "make_pptx"
    description = (
        "Build a PowerPoint (.pptx) from a list of `slides`. Each slide: `title`, `bullets` (list) or "
        "`subtitle`, an optional `image` (path, e.g. a diagram PNG), and optional `notes` (speaker "
        "notes — put the narration here). 16:9 by default. Returns the file path."
    )
    label = "Make PPTX"
    concurrency = "parallel"
    parameters = {
        "type": "object",
        "required": ["out_path", "slides"],
        "properties": {
            "out_path": {"type": "string", "description": "Output .pptx path (absolute or relative to workspace)."},
            "slides": {
                "type": "array", "minItems": 1,
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "subtitle": {"type": "string"},
                        "image": {"type": "string", "description": "Image path to place on the slide."},
                        "notes": {"type": "string", "description": "Speaker notes (e.g. the narration)."},
                    },
                },
            },
            "title_size": {"type": "integer", "description": "Title font pt. Default 36."},
            "body_size": {"type": "integer", "description": "Body font pt. Default 20."},
        },
    }

    def __init__(self, config):
        self.config = config

    def _resolve(self, p: str) -> Path:
        path = Path(p)
        if path.is_absolute():
            return path
        ws = current_workspace(str(getattr(self.config, "workspace", "."))) or "."
        return Path(ws) / path

    def _run(self, params: dict) -> dict:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except Exception as e:
            raise RuntimeError(f"python-pptx not available ({e}); install it: pip install python-pptx")

        out = self._resolve(params["out_path"])
        out.parent.mkdir(parents=True, exist_ok=True)
        title_size = int(params.get("title_size", 36))
        body_size = int(params.get("body_size", 20))

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        SW, SH = prs.slide_width, prs.slide_height

        for s in params["slides"]:
            slide = prs.slides.add_slide(blank)
            has_image = bool(s.get("image"))
            body_w = Inches(7.0) if has_image else Inches(12.13)

            if s.get("title"):
                tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), SW - Inches(1.2), Inches(1.1))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.text = s["title"]
                tf.paragraphs[0].font.size = Pt(title_size)
                tf.paragraphs[0].font.bold = True

            body = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), body_w, Inches(5.0)).text_frame
            body.word_wrap = True
            if s.get("bullets"):
                for j, b in enumerate(s["bullets"]):
                    para = body.paragraphs[0] if j == 0 else body.add_paragraph()
                    para.text = f"•  {b}"
                    para.font.size = Pt(body_size)
                    para.space_after = Pt(10)
            elif s.get("subtitle"):
                body.text = s["subtitle"]
                body.paragraphs[0].font.size = Pt(body_size + 4)

            if has_image:
                img = self._resolve(s["image"])
                if img.is_file():
                    box_l, box_t = Inches(7.7), Inches(1.8)
                    box_w, box_h = Inches(5.0), Inches(4.8)
                    pic = slide.shapes.add_picture(str(img), box_l, box_t, width=box_w)
                    if pic.height > box_h:
                        scale = box_h / pic.height
                        pic.width = int(pic.width * scale)
                        pic.height = int(pic.height * scale)
                    pic.left = box_l + (box_w - pic.width) // 2
                    pic.top = box_t + (box_h - pic.height) // 2

            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s["notes"]

        prs.save(str(out))
        return {"path": str(out), "slides": len(params["slides"])}

    async def execute(self, tool_call_id, params, abort, on_update=None):
        try:
            r = await asyncio.to_thread(self._run, params)
        except Exception as e:
            return ToolResult.text(f"make_pptx failed: {e}", is_error=True)
        return ToolResult.text(f"Wrote {r['slides']}-slide deck -> {r['path']}.", details=r)
